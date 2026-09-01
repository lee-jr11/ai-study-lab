import os
import json
import tempfile
import logging
from flask import Flask, render_template, request, jsonify
from google import genai
from google.genai import types
from dotenv import load_dotenv
from pptx import Presentation
import zipfile
from werkzeug.utils import secure_filename
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import time
import hmac
import hashlib
import secrets
from functools import wraps
from urllib.request import urlopen
import jwt  # PyJWT — verifies Firebase ID tokens server-side
from cryptography.x509 import load_pem_x509_certificate

# ──────────────────────────────────────────────
# Bootstrap
# ──────────────────────────────────────────────
load_dotenv()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024  # Hard 10 MB cap
app.config['TEMPLATES_AUTO_RELOAD'] = True  # Dev: pick up template edits without a restart

# In-memory rate limiter (no Redis needed on free tier)
limiter = Limiter(
    get_remote_address,
    app=app,
    default_limits=["200 per hour"],
    storage_uri="memory://"
)

# Gemini client
client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

# ──────────────────────────────────────────────
# Constants / whitelists
# ──────────────────────────────────────────────
ALLOWED_EXTENSIONS = {'.pdf', '.pptx'}

# File types that look like ZIP (PK magic) but are NOT PPTX — reject explicitly
# so python-pptx never sees them (it would crash on a .docx, for example).
FAKE_PPTX_EXTENSIONS = {'.docx', '.xlsx', '.zip', '.jar', '.apk', '.odt'}

# Magic-byte signatures: (expected_header, byte_length_to_check)
MAGIC_BYTES = {
    '.pdf':  (b'%PDF', 4),
    '.pptx': (b'PK',   2),   # PPTX is an Office Open XML ZIP
}

DIFFICULTY_OPTIONS = {'easy', 'medium', 'hard'}
MODE_OPTIONS       = {'quiz', 'flashcard'}
# Model strategy: Pro first, Flash fallback chain.
# The app tries PRIMARY_MODEL first, then each FALLBACK_MODEL in order until
# one responds. This handles 503 "high demand" overloads, 404 deprecations,
# and 429 quota limits automatically.
#
# NOTE: Pro models need a billing-enabled Google AI Studio account.
# On a free tier the Pro call will fail with 429 and the app will
# automatically fall back to Flash — so the site always keeps working.
PRIMARY_MODEL   = 'gemini-3.1-pro-preview'
FALLBACK_MODELS = ['gemini-2.5-flash', 'gemini-3.5-flash', 'gemini-3.1-flash-lite', 'gemini-flash-latest']

# Upload-hardening limits (defence-in-depth on top of the 10 MB upload cap)
MAX_PPTX_UNCOMPRESSED_BYTES = 100 * 1024 * 1024   # 100 MB total decompressed
MAX_EXTRACTED_TEXT_CHARS    = 200_000             # text sent to Gemini per doc


# ──────────────────────────────────────────────
# Security: Firebase ID-token verification & file ownership
# ──────────────────────────────────────────────
FIREBASE_PROJECT_ID = 'lee-ai-lab'  # must match the frontend Firebase config
GOOGLE_CERTS_URL = ('https://www.googleapis.com/robot/v1/metadata/x509/'
                    'securetoken@system.gserviceaccount.com')

_cert_cache = {'certs': None, 'fetched_at': 0.0}


def _get_google_certs():
    """Fetch & cache Google's public signing certificates (refresh hourly)."""
    if _cert_cache['certs'] is None or time.time() - _cert_cache['fetched_at'] > 3600:
        with urlopen(GOOGLE_CERTS_URL, timeout=10) as resp:
            _cert_cache['certs'] = json.loads(resp.read().decode('utf-8'))
        _cert_cache['fetched_at'] = time.time()
    return _cert_cache['certs']


class AuthError(Exception):
    """Raised when a Firebase ID token cannot be trusted."""


def verify_firebase_id_token(id_token: str) -> str:
    """Verify a Firebase ID token (signature, audience, issuer, expiry)
    and return the caller's uid."""
    try:
        header = jwt.get_unverified_header(id_token)
        kid = header.get('kid')
        if not kid:
            raise AuthError('Malformed authentication token.')
        certs = _get_google_certs()
        if kid not in certs:
            raise AuthError('Unrecognised token key. Please sign in again.')
            
        cert_str = certs[kid]
        cert_obj = load_pem_x509_certificate(cert_str.encode('utf-8'))
        public_key = cert_obj.public_key()
        
        payload = jwt.decode(
            id_token,
            public_key,
            algorithms=['RS256'],
            audience=FIREBASE_PROJECT_ID,
            issuer=f'https://securetoken.google.com/{FIREBASE_PROJECT_ID}',
            leeway=60,
        )
    except AuthError:
        raise
    except jwt.ExpiredSignatureError:
        raise AuthError('Your session expired. Please sign in again.')
    except jwt.PyJWTError as e:
        raise AuthError(f'Invalid authentication token: {str(e)}')
    uid = payload.get('user_id') or payload.get('sub')
    if not uid:
        raise AuthError('Token is missing a user identity.')
    return uid


def require_auth(f):
    """Flask decorator: verify the Bearer Firebase ID token, inject user_uid."""
    @wraps(f)
    def wrapper(*args, **kwargs):
        auth_header = request.headers.get('Authorization', '')
        if not auth_header.startswith('Bearer '):
            return jsonify({'error': 'Authentication required. Please sign in.'}), 401
        try:
            uid = verify_firebase_id_token(auth_header[7:].strip())
        except AuthError as exc:
            return jsonify({'error': str(exc)}), 401
        return f(user_uid=uid, *args, **kwargs)
    return wrapper


# Gemini Files API name -> owning Firebase uid (in-memory).
# Bound on exposure: Gemini auto-deletes hosted files after ~48h, so entries
# lost in a server restart simply expire server-side shortly afterwards.
file_owners = {}

# ── Short-lived signed "cleanup tickets" ──
# navigator.sendBeacon (used for tab-close cleanup) cannot send custom
# Authorization headers, so /init_chat issues an HMAC ticket bound to
# (uid, file_id, expiry) that /cleanup accepts as an alternative credential.
_CLEANUP_SECRET_PATH = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), '.cleanup_secret')
_cleanup_secret_cache = None


def _cleanup_secret() -> str:
    global _cleanup_secret_cache
    if _cleanup_secret_cache is None:
        if os.path.exists(_CLEANUP_SECRET_PATH):
            with open(_CLEANUP_SECRET_PATH, 'r', encoding='utf-8') as fh:
                _cleanup_secret_cache = fh.read().strip()
        else:
            generated = secrets.token_hex(32)
            with open(_CLEANUP_SECRET_PATH, 'w', encoding='utf-8') as fh:
                fh.write(generated)
            _cleanup_secret_cache = generated
    return _cleanup_secret_cache


def make_cleanup_ticket(uid: str, file_id: str, ttl_seconds: int = 7200) -> str:
    exp = int(time.time()) + ttl_seconds
    msg = f'{uid}|{file_id}|{exp}'.encode('utf-8')
    sig = hmac.new(_cleanup_secret().encode('utf-8'), msg, hashlib.sha256).hexdigest()
    return f'{exp}.{sig}'


def verify_cleanup_ticket(uid: str, file_id: str, ticket: str) -> bool:
    try:
        exp_str, sig = ticket.split('.', 1)
        exp = int(exp_str)
    except (ValueError, AttributeError):
        return False
    if exp < time.time():
        return False
    msg = f'{uid}|{file_id}|{exp}'.encode('utf-8')
    expected = hmac.new(_cleanup_secret().encode('utf-8'), msg, hashlib.sha256).hexdigest()
    return hmac.compare_digest(sig, expected)


# ──────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────
def validate_magic(file_bytes: bytes, ext: str) -> bool:
    """Return True only if file_bytes starts with the expected magic header."""
    entry = MAGIC_BYTES.get(ext)
    if entry is None:
        return False
    sig, length = entry
    return file_bytes[:length] == sig


def validate_pptx_archive(file_path: str):
    """
    Pre-parse safety check for PPTX files (which are ZIP archives).
    Reads only central-directory headers - no decompression - so a crafted
    'zip bomb' (tiny file declaring gigabytes uncompressed) is rejected
    cheaply before python-pptx ever touches it.
    Returns None when acceptable, otherwise a user-facing error message.
    """
    try:
        with zipfile.ZipFile(file_path) as zf:
            bad = zf.testzip()
            if bad is not None:
                return 'The PPTX archive appears to be corrupted.'
            total_uncompressed = sum(info.file_size for info in zf.infolist())
    except zipfile.BadZipFile:
        return 'The PPTX file is not a valid presentation archive.'
    except Exception:
        return 'The PPTX file could not be read.'
    if total_uncompressed > MAX_PPTX_UNCOMPRESSED_BYTES:
        logger.warning(
            f'pptx rejected: uncompressed size {total_uncompressed} exceeds limit'
        )
        return 'This PPTX expands to too much data to be processed safely.'
    return None


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    parts = []
    total_chars = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text)
                total_chars += len(shape.text)
                # Bound the text forwarded to Gemini (token-abuse guard)
                if total_chars >= MAX_EXTRACTED_TEXT_CHARS:
                    logger.warning('pptx text extraction hit character cap')
                    return '\n'.join(parts)[:MAX_EXTRACTED_TEXT_CHARS]
    return '\n'.join(parts)


def _strip_json_fences(text: str) -> str:
    """Remove markdown code fences Gemini sometimes wraps around JSON."""
    t = text.strip()
    if t.startswith('```'):
        newline_pos = t.find('\n')
        t = t[newline_pos + 1:] if newline_pos != -1 else ''
        if t.rstrip().endswith('```'):
            t = t.rstrip()[:-3]
    return t.strip()


def _clean_str(value) -> str:
    """Coerce a value to a trimmed string; '' when missing or not a string."""
    return value.strip() if isinstance(value, str) else ''


def validate_generated_payload(raw_text: str, mode: str, requested_count: int):
    """
    Validate & sanitize Gemini's JSON against the quiz/flashcard schema.

    Returns (items, error_message):
      - success: (non-empty list of clean dicts, None)
      - failure: (None, user-friendly error message)

    Individual malformed items are DROPPED (one bad question must not kill
    nine good ones). Only structural failures are fatal.
    """
    try:
        data = json.loads(_strip_json_fences(raw_text))
    except json.JSONDecodeError:
        return None, 'The AI returned an unexpected response. Please try again.'

    if not isinstance(data, dict) or not isinstance(data.get('quiz'), list):
        return None, 'The AI returned an unexpected response. Please try again.'

    items = []
    for entry in data['quiz']:
        if not isinstance(entry, dict):
            continue

        if mode == 'flashcard':
            term       = _clean_str(entry.get('term'))
            definition = _clean_str(entry.get('definition'))
            if term and definition:
                items.append({'term': term, 'definition': definition})
            continue

        question    = _clean_str(entry.get('question'))
        answer      = _clean_str(entry.get('answer'))
        explanation = _clean_str(entry.get('explanation'))
        raw_options = entry.get('options')

        options = []
        if isinstance(raw_options, list):
            options = [o.strip() for o in raw_options
                       if isinstance(o, str) and o.strip()]

        # Hard requirements: real question, sane option list, and an answer
        # that EXACTLY matches one of the options (else scoring breaks).
        if not question or not (2 <= len(options) <= 6):
            continue
        if not answer or answer not in options:
            continue

        item = {'question': question, 'options': options, 'answer': answer}
        if explanation:
            item['explanation'] = explanation
        items.append(item)

    if not items:
        return None, ('The AI could not produce valid content from this document. '
                      'Please try again, or use a different document.')

    if len(items) > requested_count:
        items = items[:requested_count]

    if len(items) < requested_count:
        logger.warning(
            f"generate_quiz: Gemini returned {len(items)} valid {mode} item(s); "
            f"{requested_count} were requested."
        )

    return items, None


# ──────────────────────────────────────────────
# Error handlers
# ──────────────────────────────────────────────
@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Maximum allowed size is 10 MB.'}), 413


@app.errorhandler(429)
def rate_limited(e):
    return jsonify({'error': 'Too many requests. Please slow down and try again shortly.'}), 429


# ──────────────────────────────────────────────
# Routes
# ──────────────────────────────────────────────
@app.route('/')
def index():
    return render_template('index.html')


@app.route('/generate', methods=['POST'])
@limiter.limit("10 per minute")
@require_auth
def generate_quiz(user_uid):
    if 'document' not in request.files:
        return jsonify({'error': 'No document uploaded.'}), 400

    file = request.files['document']
    if not file.filename:
        return jsonify({'error': 'No file selected.'}), 400

    # Sanitise filename & check extension
    filename = secure_filename(file.filename)
    file_ext = os.path.splitext(filename)[1].lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        if file_ext in FAKE_PPTX_EXTENSIONS:
            return jsonify({'error': f'{file_ext.upper().lstrip(".") } files are not supported. Please convert your document to PDF or PPTX first.'}), 400
        return jsonify({'error': 'Only PDF and PPTX files are accepted.'}), 400

    # Whitelist & clamp inputs
    difficulty = request.form.get('difficulty', 'medium')
    if difficulty not in DIFFICULTY_OPTIONS:
        difficulty = 'medium'

    try:
        num_questions = max(1, min(30, int(request.form.get('num_questions', 10))))
    except (ValueError, TypeError):
        num_questions = 10

    mode = request.form.get('mode', 'quiz')
    if mode not in MODE_OPTIONS:
        mode = 'quiz'

    # Read bytes & validate magic header
    file_bytes = file.read()
    if not validate_magic(file_bytes, file_ext):
        return jsonify({'error': 'File content does not match its type. Please re-export and try again.'}), 400

    temp_file_path = None
    gemini_file = None

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tf:
            tf.write(file_bytes)
            temp_file_path = tf.name

        if mode == 'flashcard':
            prompt = (
                f"You are an expert tutor. Create exactly {num_questions} flashcards from the provided document.\n"
                f"Difficulty: {difficulty}.\n\n"
                "Respond ONLY with a valid JSON object — no markdown, no extra text.\n"
                "The object must have a single key \"quiz\" containing an array of objects.\n"
                "Each object must have exactly:\n"
                "- \"term\": the key concept or vocabulary word\n"
                "- \"definition\": a clear, concise explanation"
            )
        else:
            prompt = (
                f"You are an expert university professor. Generate exactly {num_questions} multiple-choice questions from the provided document.\n"
                f"Difficulty: {difficulty}.\n\n"
                "Respond ONLY with a valid JSON object — no markdown, no extra text.\n"
                "The object must have a single key \"quiz\" containing an array of objects.\n"
                "Each object must have exactly:\n"
                "- \"question\": the question text\n"
                "- \"options\": array of exactly 4 possible answers\n"
                "- \"answer\": the correct option string (must exactly match one of the options)\n"
                "- \"explanation\": 1-2 sentence explanation of why this answer is correct"
            )

        contents = [prompt]

        if file_ext == '.pptx':
            zip_error = validate_pptx_archive(temp_file_path)
            if zip_error:
                return jsonify({'error': zip_error}), 400
            pptx_text = extract_text_from_pptx(temp_file_path)
            if not pptx_text.strip():
                return jsonify({'error': 'The PPTX file appears to have no readable text content.'}), 400
            contents.append(f'Document content:\n\n{pptx_text}')
        else:
            gemini_file = client.files.upload(file=temp_file_path)
            contents.append(gemini_file)

        # Generate + validate, with ONE automatic retry when the AI's output
        # is unusable (unparseable, wrong shape, or zero valid items).
        parsed_items = None
        last_error   = 'The AI returned an unexpected response. Please try again.'
        for attempt_no in (1, 2):
            try:
                # Try models in priority order: Pro first, then each Flash fallback.
                # Falls through on ANY failure (503 overload, 404 not found, timeout, etc.)
                response = None
                last_api_error = None
                for model_name in [PRIMARY_MODEL] + FALLBACK_MODELS:
                    try:
                        response = client.models.generate_content(
                            model=model_name,
                            contents=contents,
                            config=types.GenerateContentConfig(response_mime_type='application/json')
                        )
                        logger.info(f'generate_quiz: {model_name} responded successfully')
                        break  # success — stop trying fallbacks
                    except Exception as api_err:
                        last_api_error = api_err
                        err_str = str(api_err).lower()
                        logger.warning(f'generate_quiz: {model_name} failed ({err_str[:120]}), trying next model...')
                if response is None:
                    raise last_api_error if last_api_error else RuntimeError('All Gemini models failed.')

                parsed_items, last_error = validate_generated_payload(
                    response.text or '', mode, num_questions
                )
            except Exception as gen_err:
                logger.error(f'generate_quiz: attempt {attempt_no} failed: {gen_err}')
                last_error   = 'Failed to generate content. The document may be too complex or large.'
                parsed_items = None

            if parsed_items is not None:
                break
            if attempt_no == 1:
                logger.warning('generate_quiz: unusable AI output — retrying once.')

        if parsed_items is None:
            return jsonify({'error': last_error}), 500

        return jsonify({'quiz': parsed_items})

    except Exception as e:
        logger.error(f'generate_quiz error: {e}')
        return jsonify({'error': f'Failed to generate content: {str(e)}'}), 500
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        if gemini_file:
            try:
                client.files.delete(name=gemini_file.name)
            except Exception:
                pass


@app.route('/init_chat', methods=['POST'])
@limiter.limit("10 per minute")
@require_auth
def init_chat(user_uid):
    if 'document' not in request.files:
        return jsonify({'error': 'No document uploaded.'}), 400

    file = request.files['document']
    if not file.filename:
        return jsonify({'error': 'No file selected.'}), 400

    filename = secure_filename(file.filename)
    file_ext = os.path.splitext(filename)[1].lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        if file_ext in FAKE_PPTX_EXTENSIONS or file_ext == '.docx':
            return jsonify({'error': f'{file_ext.upper().lstrip(".") } files are not supported. Please convert your document to PDF or PPTX first.'}), 400
        return jsonify({'error': 'Only PDF and PPTX files are accepted.'}), 400

    file_bytes = file.read()
    if not validate_magic(file_bytes, file_ext):
        return jsonify({'error': 'File content does not match its type.'}), 400

    temp_file_path = None
    txt_path = None
    gemini_file = None

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tf:
            tf.write(file_bytes)
            temp_file_path = tf.name

        if file_ext == '.pptx':
            zip_error = validate_pptx_archive(temp_file_path)
            if zip_error:
                return jsonify({'error': zip_error}), 400
            pptx_text = extract_text_from_pptx(temp_file_path)
            # Consistency with /generate: never start a chat session on a
            # text-less deck (no slides, or slides whose shapes are all empty).
            if not pptx_text.strip():
                return jsonify({'error': 'The PPTX file appears to have no readable text content.'}), 400
            with tempfile.NamedTemporaryFile(delete=False, suffix='.txt', mode='w', encoding='utf-8') as txt_f:
                txt_f.write(pptx_text)
                txt_path = txt_f.name
            gemini_file = client.files.upload(file=txt_path)
        else:
            gemini_file = client.files.upload(file=temp_file_path)

        # Bind this hosted document to its owner and issue a short-lived
        # signed cleanup ticket (used by the tab-close beacon, which cannot
        # send Authorization headers).
        file_owners[gemini_file.name] = user_uid
        return jsonify({
            'file_id': gemini_file.name,
            'ct': make_cleanup_ticket(user_uid, gemini_file.name),
        })

    except Exception as e:
        logger.error(f'init_chat error: {e}')
        if gemini_file:
            try:
                client.files.delete(name=gemini_file.name)
            except Exception:
                pass
        return jsonify({'error': 'Failed to process the document. Please try again.'}), 500
    finally:
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        if txt_path and os.path.exists(txt_path):
            os.remove(txt_path)


@app.route('/chat', methods=['POST'])
@limiter.limit("30 per minute")
@require_auth
def chat_with_document(user_uid):
    file_id      = request.form.get('file_id', '').strip()
    user_message = request.form.get('message', '').strip()
    history_raw  = request.form.get('history', '[]')

    if not file_id:
        return jsonify({'error': 'Missing file context. Please restart the session.'}), 400
    if not user_message:
        return jsonify({'error': 'Message cannot be empty.'}), 400
    if len(user_message) > 2000:
        return jsonify({'error': 'Message too long. Please keep it under 2000 characters.'}), 400

    # Ownership check: a valid session may only query documents the caller
    # uploaded themselves.
    if file_owners.get(file_id) != user_uid:
        logger.warning(f'chat: uid {user_uid} denied access to foreign file {file_id}')
        return jsonify({'error': 'You do not have access to this document session.'}), 403

    try:
        history_list = json.loads(history_raw)
        if not isinstance(history_list, list):
            history_list = []
    except (json.JSONDecodeError, ValueError):
        history_list = []

    # Cap to last 20 messages to control token usage on free tier
    history_list = history_list[-20:]

    # Defensive dedupe: if a client already included the current message at
    # the end of its history, strip it so the question is never sent twice.
    if history_list and isinstance(history_list[-1], dict):
        last_entry = history_list[-1]
        if (last_entry.get('role') == 'user'
                and str(last_entry.get('content', '')).strip() == user_message):
            history_list = history_list[:-1]

    try:
        prompt = (
            "You are an expert, friendly AI tutor. "
            "Use the attached document to answer the student's question.\n\n"
        )
        if history_list:
            prompt += "Recent conversation:\n"
            for msg in history_list:
                if isinstance(msg, dict) and 'role' in msg and 'content' in msg:
                    role    = 'Student' if msg['role'] == 'user' else 'Tutor'
                    content = str(msg['content'])[:500]
                    prompt += f"{role}: {content}\n"

        prompt += f"\nStudent: {user_message}\nTutor:"

        try:
            gemini_file = client.files.get(name=file_id)
        except Exception as file_err:
            # The server restarted (wiping file_owners) or the Gemini file
            # expired (~48 h TTL). Either way, the session is unrecoverable
            # — tell the user clearly so they know to re-upload.
            logger.warning(f'chat: gemini file {file_id} not found: {file_err}')
            # Clean up stale owner record if it somehow survived
            file_owners.pop(file_id, None)
            return jsonify({
                'error': 'Your document session has expired. Please go back and re-upload your file to start a new chat.'
            }), 410

        # Try models in priority order: Pro first, then each Flash fallback.
        response = None
        last_api_error = None
        for model_name in [PRIMARY_MODEL] + FALLBACK_MODELS:
            try:
                response = client.models.generate_content(
                    model=model_name,
                    contents=[gemini_file, prompt]
                )
                logger.info(f'chat: {model_name} responded successfully')
                break  # success — stop trying fallbacks
            except Exception as api_err:
                last_api_error = api_err
                err_str = str(api_err).lower()
                logger.warning(f'chat: {model_name} failed ({err_str[:120]}), trying next model...')
        if response is None:
            raise last_api_error if last_api_error else RuntimeError('All Gemini models failed.')

        return jsonify({'reply': response.text})

    except Exception as e:
        logger.error(f'chat error: {e}')
        return jsonify({'error': 'Something went wrong. Please try again.'}), 500


@app.route('/models')
def list_models():
    try:
        models = [m.name for m in client.models.list()]
        return jsonify(models)
    except Exception as e:
        return jsonify({'error': str(e)})

@app.route('/cleanup', methods=['POST'])
@limiter.limit("60 per minute")
def cleanup_chat():
    """
    Called when a chat session ends (Back button, sign-out, tab close).
    Authorised either by a Bearer Firebase ID token OR by the short-lived
    signed cleanup ticket issued at /init_chat (needed because the
    tab-close beacon cannot send custom Authorization headers).
    Unauthorised or unowned requests are silent no-ops so attackers can't
    probe for or delete other users' hosted documents.
    """
    file_id = request.form.get('file_id', '').strip()

    uid = None
    auth_header = request.headers.get('Authorization', '')
    if auth_header.startswith('Bearer '):
        try:
            uid = verify_firebase_id_token(auth_header[7:].strip())
        except AuthError:
            uid = None
    if uid is None and file_id:
        ticket = request.form.get('ct', '').strip()
        owner = file_owners.get(file_id)
        if owner and ticket and verify_cleanup_ticket(owner, file_id, ticket):
            uid = owner

    if file_id and uid is not None:
        owner = file_owners.get(file_id)
        if owner == uid:
            file_owners.pop(file_id, None)
            try:
                client.files.delete(name=file_id)
                logger.info(f'cleanup: deleted Gemini file {file_id}')
            except Exception as e:
                logger.warning(f'cleanup: could not delete {file_id}: {e}')
        else:
            logger.warning(f'cleanup: ignored unowned/foreign file {file_id}')
    return jsonify({'status': 'ok'})


# ──────────────────────────────────────────────
# Entry point
# ──────────────────────────────────────────────
if __name__ == '__main__':
    # SECURITY: Flask's interactive debugger allows arbitrary remote code
    # execution if ever exposed, so it can no longer be enabled through an
    # environment variable that might accidentally carry into production.
    # Day-to-day dev convenience is covered by TEMPLATES_AUTO_RELOAD instead.
    if os.getenv('FLASK_DEBUG', '').strip().lower() == 'true':
        logger.warning(
            'FLASK_DEBUG=true detected but IGNORED - the interactive debugger '
            'is permanently disabled for security reasons.'
        )
    app.run(debug=False)

